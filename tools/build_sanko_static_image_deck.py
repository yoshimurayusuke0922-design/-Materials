from __future__ import annotations

from pathlib import Path
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
IMAGE_DIR = ROOT / "out" / "sanko_html" / "slide_images"
OUT_DIR = ROOT / "out" / "sanko_static_slides"
OUT_PPTX = OUT_DIR / "sanko_kikan_static_image_deck.pptx"
OUT_PDF = OUT_DIR / "sanko_kikan_static_image_deck.pdf"
OUT_CONTACT_SHEET = OUT_DIR / "sanko_kikan_static_image_deck_contact_sheet.png"

SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
THUMB_W = 420
THUMB_H = 236


def slide_images() -> list[Path]:
    paths = sorted(IMAGE_DIR.glob("*.png"))
    if not paths:
        raise RuntimeError(f"No PNG slide images found: {IMAGE_DIR}")
    return paths


def picture_fit(path: Path) -> tuple[float, float, float, float]:
    with Image.open(path) as image:
        img_w, img_h = image.size
    slide_ratio = SLIDE_W_IN / SLIDE_H_IN
    image_ratio = img_w / img_h
    if image_ratio >= slide_ratio:
        width = SLIDE_W_IN
        height = SLIDE_W_IN / image_ratio
    else:
        height = SLIDE_H_IN
        width = SLIDE_H_IN * image_ratio
    left = (SLIDE_W_IN - width) / 2
    top = (SLIDE_H_IN - height) / 2
    return left, top, width, height


def write_pptx(paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for path in paths:
        slide = prs.slides.add_slide(blank)
        left, top, width, height = picture_fit(path)
        slide.shapes.add_picture(
            str(path),
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height),
        )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    with zipfile.ZipFile(OUT_PPTX) as z:
        bad = z.testzip()
        if bad:
            raise RuntimeError(f"Broken PPTX entry: {bad}")


def write_pdf(paths: list[Path]) -> None:
    pages = [Image.open(path).convert("RGB") for path in paths]
    try:
        pages[0].save(OUT_PDF, save_all=True, append_images=pages[1:], resolution=150.0)
    finally:
        for page in pages:
            page.close()


def font(size: int) -> ImageFont.ImageFont:
    for candidate in [
        r"C:\Windows\Fonts\YuGothB.ttc",
        r"C:\Windows\Fonts\meiryob.ttc",
        r"C:\Windows\Fonts\arial.ttf",
    ]:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def write_contact_sheet(paths: list[Path]) -> None:
    thumbs: list[Image.Image] = []
    for path in paths:
        image = Image.open(path).convert("RGB")
        image.thumbnail((THUMB_W, THUMB_H), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (THUMB_W, THUMB_H), "white")
        canvas.paste(image, ((THUMB_W - image.width) // 2, (THUMB_H - image.height) // 2))
        thumbs.append(canvas)

    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    pad = 28
    label_h = 34
    sheet = Image.new(
        "RGB",
        (cols * THUMB_W + (cols + 1) * pad, rows * (THUMB_H + label_h) + (rows + 1) * pad),
        (245, 247, 250),
    )
    draw = ImageDraw.Draw(sheet)
    label_font = font(22)
    for i, thumb in enumerate(thumbs):
        row, col = divmod(i, cols)
        x = pad + col * (THUMB_W + pad)
        y = pad + row * (THUMB_H + label_h + pad)
        sheet.paste(thumb, (x, y))
        draw.text((x + 8, y + THUMB_H + 7), f"Slide {i + 1:02d}", fill=(24, 24, 24), font=label_font)
    sheet.save(OUT_CONTACT_SHEET, quality=92)
    for thumb in thumbs:
        thumb.close()


def main() -> None:
    paths = slide_images()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    write_pptx(paths)
    write_pdf(paths)
    write_contact_sheet(paths)
    print(f"slides={len(paths)}")
    print(OUT_PPTX)
    print(OUT_PDF)
    print(OUT_CONTACT_SHEET)


if __name__ == "__main__":
    main()
