from __future__ import annotations

from pathlib import Path
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
PROJECT = ROOT / "projects" / "2026-04-24-period-management-proposal"
IMAGE_DIR = PROJECT / "image2_only"
OUT = PROJECT / "株式会社三幸_基幹管理システムご提案資料_image2_only.pptx"
OVERVIEW = IMAGE_DIR / "overview.png"

SLIDE_W = 13.333333
SLIDE_H = 7.5


def image_files() -> list[Path]:
    return sorted(IMAGE_DIR.glob("slide_*.png"))


def write_pptx(paths: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    for path in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    with zipfile.ZipFile(OUT) as z:
        bad = z.testzip()
        if bad:
            raise RuntimeError(f"Broken pptx entry: {bad}")


def font(size: int) -> ImageFont.ImageFont:
    for candidate in [
        r"C:\Windows\Fonts\YuGothB.ttc",
        r"C:\Windows\Fonts\meiryob.ttc",
        r"C:\Windows\Fonts\arial.ttf",
    ]:
        p = Path(candidate)
        if p.exists():
            return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def write_overview(paths: list[Path]) -> None:
    thumbs = []
    for path in paths:
        img = Image.open(path).convert("RGB")
        img.thumbnail((420, 236), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (420, 236), "white")
        canvas.paste(img, ((420 - img.width) // 2, (236 - img.height) // 2))
        thumbs.append(canvas)

    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    pad = 28
    label_h = 34
    out = Image.new("RGB", (cols * 420 + (cols + 1) * pad, rows * (236 + label_h) + (rows + 1) * pad), (245, 245, 245))
    draw = ImageDraw.Draw(out)
    fnt = font(22)
    for i, thumb in enumerate(thumbs):
        row, col = divmod(i, cols)
        x = pad + col * (420 + pad)
        y = pad + row * (236 + label_h + pad)
        out.paste(thumb, (x, y))
        draw.text((x + 8, y + 242), f"Slide {i + 1:02d}", fill=(30, 30, 30), font=fnt)
    out.save(OVERVIEW, quality=92)


def main() -> None:
    paths = image_files()
    if not paths:
        raise RuntimeError(f"No slide images found in {IMAGE_DIR}")
    write_pptx(paths)
    write_overview(paths)
    print(OUT)
    print(OVERVIEW)


if __name__ == "__main__":
    main()
