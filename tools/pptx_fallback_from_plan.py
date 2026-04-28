from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt


def rgb(hex_color: str | None) -> RGBColor:
    value = (hex_color or "#000000").replace("#", "")
    if value == "transparent":
        value = "FFFFFF"
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    value = (value + "000000")[:6]
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def pt(value: float | int):
    return Pt(float(value or 0))


def shape_type(name: str):
    return {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "roundRectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "rightArrow": MSO_SHAPE.RIGHT_ARROW,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    }.get(name, MSO_SHAPE.RECTANGLE)


def add_shape(slide, element):
    shp = slide.shapes.add_shape(
        shape_type(element.get("shape")),
        pt(element["x"]),
        pt(element["y"]),
        pt(element["width"]),
        pt(element["height"]),
    )
    fill = element.get("fill")
    if fill and fill != "transparent":
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    border = element.get("border")
    if border and border != "transparent":
        shp.line.color.rgb = rgb(border)
        shp.line.width = pt(element.get("border_width", 1))
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, element):
    box = slide.shapes.add_textbox(
        pt(element["x"]),
        pt(element["y"]),
        pt(element["width"]),
        pt(element["height"]),
    )
    if element.get("rotation"):
        box.rotation = element["rotation"]
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(1)
    frame.margin_right = Pt(1)
    frame.margin_top = Pt(1)
    frame.margin_bottom = Pt(1)
    valign = element.get("valign")
    if valign in ("middle", "center"):
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        frame.vertical_anchor = MSO_ANCHOR.TOP
    lines = str(element.get("text", "")).split("\n") or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = {
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "left": PP_ALIGN.LEFT,
        }.get(element.get("align"), PP_ALIGN.LEFT)
        for run in paragraph.runs:
            run.font.name = element.get("font_family", "Yu Gothic")
            run.font.size = pt(element.get("font_size", 12))
            run.font.bold = bool(element.get("bold"))
            run.font.color.rgb = rgb(element.get("color", "#101214"))
    return box


def add_line(slide, element):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        pt(element["x1"]),
        pt(element["y1"]),
        pt(element["x2"]),
        pt(element["y2"]),
    )
    conn.line.color.rgb = rgb(element.get("color", "#f80612"))
    conn.line.width = pt(element.get("width", 1))
    conn.line.dash_style = MSO_LINE.SOLID
    return conn


def add_image(slide, element, asset_base: Path):
    src = element.get("src")
    if not src:
        return None
    path = Path(src)
    if not path.is_absolute():
        path = asset_base / src
    if not path.exists():
        return None
    return slide.shapes.add_picture(
        str(path),
        pt(element["x"]),
        pt(element["y"]),
        width=pt(element["width"]),
        height=pt(element["height"]),
    )


def build(plan_path: Path, asset_base: Path, output_path: Path) -> None:
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = pt(plan["page"]["width"])
    prs.slide_height = pt(plan["page"]["height"])
    blank = prs.slide_layouts[6]
    while len(prs.slides) > 0:
        break
    for slide_plan in plan["slides"]:
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(slide_plan.get("background", "#ffffff"))
        for element in slide_plan.get("elements", []):
            kind = element.get("type")
            if kind == "shape":
                add_shape(slide, element)
            elif kind == "text":
                add_text(slide, element)
            elif kind == "line":
                add_line(slide, element)
            elif kind == "image":
                add_image(slide, element, asset_base)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--plan", default="intermediate/google_slides_plan.json")
    parser.add_argument("--asset-base", default="out/sanko_html/rebuilt_slides")
    parser.add_argument("--out", default="output/fallback_editable_from_plan.pptx")
    args = parser.parse_args()
    build(Path(args.plan), Path(args.asset_base), Path(args.out))
    print(args.out)


if __name__ == "__main__":
    main()

