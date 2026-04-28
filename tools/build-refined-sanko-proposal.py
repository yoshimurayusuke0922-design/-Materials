from __future__ import annotations

from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PROJECT = ROOT / "projects" / "2026-04-24-period-management-proposal"
REF_OPENING = ROOT / "reference_slides" / "library" / "google-slides-1jynt98j" / "slides"
REF_MOCK = ROOT / "reference_slides" / "library" / "period-management-mock" / "slides"
OUT = PROJECT / "株式会社三幸_基幹管理システムご提案資料_refined.pptx"
DESIGN_CONCEPTS = PROJECT / "design_concepts"
MOCK_DASHBOARD_IMAGE = DESIGN_CONCEPTS / "mock_ui_dashboard_crop.png"

RED = RGBColor(255, 24, 24)
RED_DARK = RGBColor(220, 0, 0)
PINK = RGBColor(255, 120, 120)
BLACK = RGBColor(18, 18, 18)
TEXT = RGBColor(37, 43, 54)
MUTED = RGBColor(100, 100, 100)
GRAY_BG = RGBColor(242, 242, 242)
GRAY_CARD = RGBColor(248, 248, 248)
LINE = RGBColor(220, 220, 220)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(23, 25, 31)

SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Yu Gothic"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_fill(shape, color: RGBColor | None) -> None:
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None = None, width: float = 0.75) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None, radius=False):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill)
    set_line(shp, line)
    return shp


def add_rect(slide, x, y, w, h, fill=None, line=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)


def add_round(slide, x, y, w, h, fill=None, line=None):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def add_oval(slide, x, y, w, h, fill=None, line=None):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill, line)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    value: str,
    size: float = 18,
    color: RGBColor = BLACK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    lines = value.split("\n") if value else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    set_fill(box, None)
    set_line(box, None)
    return box


def add_bullets(slide, x, y, w, h, items: list[str], size=15, color=TEXT, bold=False):
    text = "\n".join(f"・{item}" for item in items)
    return add_text(slide, x, y, w, h, text, size, color, bold)


def prepare_assets() -> None:
    """Create reusable image2-derived figure crops with no final Japanese text."""
    if MOCK_DASHBOARD_IMAGE.exists():
        return
    from PIL import Image

    src = DESIGN_CONCEPTS / "concept_03_mock_ui.png"
    DESIGN_CONCEPTS.mkdir(parents=True, exist_ok=True)
    img = Image.open(src).convert("RGB")
    crop = img.crop((18, 105, 1320, 868))
    crop.save(MOCK_DASHBOARD_IMAGE, quality=95)


def add_page_rail(slide, page: int) -> None:
    add_rect(slide, 0, 0, 0.44, SLIDE_H, RED, RED)
    add_text(slide, 0.08, 6.90, 0.28, 0.38, str(page), 18, WHITE, True, PP_ALIGN.CENTER)


def add_header(slide, page: int, chapter: str, title: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_BG
    add_page_rail(slide, page)
    add_text(slide, 0.74, 0.34, 1.3, 0.28, chapter, 16, RED, True)
    add_text(slide, 0.74, 0.70, 7.8, 0.48, title, 24, RGBColor(78, 78, 78), True)


def card(slide, x, y, w, h, title: str, body: str, title_size=15, body_size=12.5):
    add_round(slide, x, y, w, h, WHITE, LINE)
    add_rect(slide, x + 0.18, y + 0.17, 0.04, 0.36, RED, RED)
    add_text(slide, x + 0.28, y + 0.13, w - 0.45, 0.38, title, title_size, RED, True)
    add_text(slide, x + 0.28, y + 0.66, w - 0.55, h - 0.78, body, body_size, TEXT)


def cover(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_BG
    add_oval(slide, 8.05, -0.10, 5.60, 7.90, RED, RED)
    add_oval(slide, -1.55, 5.75, 4.15, 2.60, RED, RED)
    add_oval(slide, 3.95, -2.20, 3.40, 3.90, rgb("ffb0b0"), rgb("ffb0b0"))
    add_rect(slide, 5.35, 0.00, 1.95, 1.85, GRAY_BG, GRAY_BG)
    add_text(slide, 0.84, 2.28, 3.30, 0.32, "株式会社三幸　御中", 15, BLACK, True)
    add_text(slide, 0.84, 3.05, 7.30, 0.72, "基幹管理システムの", 34, BLACK, True)
    add_text(slide, 0.84, 3.92, 6.20, 0.72, "ご提案資料", 34, BLACK, True)
    add_text(slide, 0.84, 5.25, 2.40, 0.32, "2026年04月25日", 16, BLACK)
    add_oval(slide, 6.10, 6.64, 0.10, 0.10, RED, RED)
    add_oval(slide, 6.28, 6.54, 0.10, 0.10, RED, RED)
    add_oval(slide, 6.28, 6.74, 0.10, 0.10, RED, RED)
    add_oval(slide, 6.46, 6.64, 0.10, 0.10, RED, RED)
    add_text(slide, 6.58, 6.43, 1.40, 0.40, "Field X", 24, BLACK, True)


def agenda(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_BG
    add_page_rail(slide, 2)
    add_text(slide, 0.74, 0.36, 1.5, 0.28, "AGENDA", 18, RED, True)
    add_text(slide, 0.74, 0.78, 1.4, 0.42, "目次", 24, RGBColor(90, 90, 90), True)
    rows = [
        ("01", "これまでにお伺いしたこと"),
        ("02", "ご提案の要旨・デモンストレーション"),
        ("03", "今後の進め方のイメージ"),
    ]
    y = 2.26
    for num, label in rows:
        add_text(slide, 1.86, y, 0.82, 0.55, num, 36, RED, True, PP_ALIGN.RIGHT)
        add_rect(slide, 2.90, y - 0.02, 0.03, 0.68, RED, RED)
        add_text(slide, 3.22, y + 0.08, 6.8, 0.48, label, 27, RED, True)
        y += 1.25


def chapter(prs, page: int, num: str, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RED
    add_text(slide, 1.15, 2.50, 3.2, 1.7, num, 118, PINK, True)
    add_text(slide, 1.78, 3.58, 6.0, 0.6, title, 26, WHITE, True)


def position_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 4, "Chapter1", "本資料の位置づけ")
    add_round(slide, 0.78, 1.48, 5.7, 3.05, WHITE, LINE)
    add_round(slide, 6.86, 1.48, 5.7, 3.05, WHITE, LINE)
    add_text(slide, 1.08, 1.76, 2.0, 0.36, "今回の目的", 17, RED, True)
    add_text(slide, 7.16, 1.76, 2.6, 0.36, "確認いただく観点", 17, RED, True)
    add_bullets(slide, 1.10, 2.28, 4.9, 1.7, [
        "確認済みモックを、提案資料として整理",
        "開発範囲・主要機能・導入後の流れを明確化",
        "画面ごとに要件と機能を確認できる形にする",
    ], 14)
    add_bullets(slide, 7.18, 2.28, 4.9, 1.7, [
        "業務に合わせたカスタム開発の方向性",
        "初期開発で優先する画面・帳票・ワークフロー",
        "将来的なAI活用まで見据えた拡張方針",
    ], 14)
    items = [
        ("提出物", "改善版モック\n画面別要件資料\n開発の進め方"),
        ("資料の方向性", "画面を軸に整理\n機能範囲を明確化\n導入後の流れを説明"),
        ("次に確認したいこと", "優先機能\n権限・帳票\n既存データ移行"),
    ]
    for i, (h, b) in enumerate(items):
        x = 0.90 + i * 4.12
        card(slide, x, 5.02, 3.45, 1.18, h, b, 15, 12)


def heard_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 5, "Chapter1", "これまでにお伺いしたこと")
    cards = [
        ("基幹システム切り替え", "既存システムのサポート終了を見据え、今年中を目途に切り替え検討が必要。"),
        ("現場定着への配慮", "IT活用への慣れに差があるため、日常業務に自然に入る画面設計が重要。"),
        ("日常業務の負荷", "退去修繕対応、契約書作成、請求・インボイス対応などの工数が大きい。"),
        ("データ活用の余地", "顧客・物件・契約・対応履歴を横断して見られる状態を作り、判断の質を高める。"),
    ]
    for i, (h, b) in enumerate(cards):
        card(slide, 0.72 + i * 3.15, 1.62, 2.72, 2.82, h, b, 14.5, 13)
    add_round(slide, 0.78, 5.15, 11.65, 0.76, rgb("fff1f1"), RED)
    add_text(slide, 1.05, 5.38, 11.1, 0.28, "本提案では、単なる画面開発ではなく「基幹業務の整理」と「今後AIを載せられる土台づくり」を同時に進めます。", 16, BLACK, True)


def proposal_summary(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 7, "Chapter2", "ご提案の要旨・デモンストレーション")
    add_round(slide, 0.72, 1.38, 11.95, 1.05, WHITE, None)
    add_text(slide, 1.08, 1.72, 1.7, 0.30, "実現すること", 18, RED, True)
    add_round(slide, 3.22, 1.53, 7.9, 0.58, RED, RED)
    add_text(slide, 3.45, 1.66, 7.45, 0.30, "基幹管理・期日管理・対応履歴を一体化し、業務全体の効率化を実現", 15.5, WHITE, True, PP_ALIGN.CENTER)
    add_round(slide, 0.95, 2.70, 11.5, 2.55, WHITE, None)
    entries = [
        ("基幹業務の再設計・標準化", ["業務フローを整理", "退去修繕・契約業務を標準化"]),
        ("物件・反響対応の効率化", ["出稿・更新状況を管理", "反響から返信までを一元化"]),
        ("契約・請求業務の自動化", ["契約書・請求書作成を省力化", "確認漏れや作業負荷を削減"]),
        ("現場対応の即時化", ["物件・顧客・履歴を参照", "外出先でも確認対応を完結"]),
    ]
    for i, (h, items) in enumerate(entries):
        x = 1.35 + (i % 2) * 5.75
        y = 3.05 + (i // 2) * 1.02
        add_text(slide, x, y, 4.2, 0.30, h, 15.5, RED, True)
        add_rect(slide, x, y + 0.33, 3.7, 0.02, RED, RED)
        add_bullets(slide, x, y + 0.48, 4.75, 0.52, items, 11.5)
    add_round(slide, 0.95, 5.58, 11.5, 0.72, WHITE, None)
    add_text(slide, 4.55, 5.86, 6.8, 0.28, "改善版モックで、画面ごとの要件と機能を確認できる形にします。", 14, BLACK, True)
    add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, 4.08, 5.80, 0.32, 0.32, RED, RED)


def app_shell(slide, x, y, w, h, title: str, active: str = "総合"):
    add_round(slide, x, y, w, h, WHITE, LINE)
    side_w = 1.18
    add_rect(slide, x, y, side_w, h, DARK, DARK)
    add_oval(slide, x + 0.18, y + 0.18, 0.22, 0.22, RED, RED)
    add_text(slide, x + 0.46, y + 0.16, 0.54, 0.22, "Field X", 10, WHITE, True)
    add_rect(slide, x + side_w, y, w - side_w, 0.44, rgb("fbfbfb"), LINE)
    add_text(slide, x + side_w + 0.20, y + 0.14, 4.0, 0.20, title, 10.5, BLACK, True)
    nav = ["総合", "物件", "契約", "修繕", "書類", "検索"]
    for i, item in enumerate(nav):
        ny = y + 0.72 + i * 0.48
        fill = RED if item == active else rgb("2f323b")
        add_round(slide, x + 0.16, ny, 0.86, 0.30, fill, fill)
        add_text(slide, x + 0.23, ny + 0.08, 0.70, 0.12, item, 6.2, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, x + 0.18, y + h - 0.55, 0.76, 0.01, rgb("51545b"), rgb("51545b"))
    add_text(slide, x + 0.25, y + h - 0.38, 0.50, 0.12, "設定", 6.0, rgb("d4d4d4"))
    return x + side_w + 0.24, y + 0.72, w - side_w - 0.46, h - 0.95


def dashboard_body(slide, x, y, w, h):
    stats = [("本日期限", "12"), ("要確認", "7"), ("遅延リスク", "3"), ("完了", "28"), ("未割当", "4")]
    bw = (w - 0.32) / 5
    for i, (label, value) in enumerate(stats):
        bx = x + i * (bw + 0.08)
        add_round(slide, bx, y, bw, 0.62, WHITE, LINE)
        add_text(slide, bx + 0.09, y + 0.11, bw - 0.18, 0.12, label, 6.5, MUTED, True)
        add_text(slide, bx + 0.11, y + 0.30, 0.36, 0.22, value, 15, RED if label == "遅延リスク" else BLACK, True)
    cols = ["退去精算", "修繕対応", "契約更新", "請求確認", "入居対応"]
    cw = (w - 0.36) / 5
    for i, col in enumerate(cols):
        bx = x + i * (cw + 0.09)
        add_round(slide, bx, y + 0.82, cw, 2.15, rgb("f8f8f8"), LINE)
        add_text(slide, bx + 0.10, y + 0.98, cw - 0.25, 0.13, col, 7.5, BLACK, True)
        for j in range(3):
            yy = y + 1.28 + j * 0.55
            add_round(slide, bx + 0.08, yy, cw - 0.16, 0.40, WHITE, LINE)
            add_rect(slide, bx + 0.08, yy, 0.03, 0.40, RED if j == 0 else rgb("bfbfbf"), None)
            add_text(slide, bx + 0.17, yy + 0.08, cw - 0.35, 0.10, f"案件 {i + 1}-{j + 1}", 5.8, BLACK, True)
            add_text(slide, bx + 0.17, yy + 0.23, cw - 0.35, 0.10, "期限 / 担当 / 状態", 4.8, MUTED)
    add_round(slide, x, y + 3.18, w * 0.56, 0.70, WHITE, LINE)
    add_text(slide, x + 0.14, y + 3.33, 1.2, 0.12, "直近対応履歴", 6.8, BLACK, True)
    for i, label in enumerate(["退去立会い", "見積確認", "契約更新"]):
        yy = y + 3.58 + i * 0.16
        add_oval(slide, x + 0.15, yy, 0.07, 0.07, RED if i == 0 else rgb("aaaaaa"), None)
        add_text(slide, x + 0.30, yy - 0.02, 1.0, 0.08, label, 4.6, MUTED)
    add_round(slide, x + w * 0.60, y + 3.18, w * 0.40, 0.70, WHITE, LINE)
    add_oval(slide, x + w * 0.65, y + 3.37, 0.34, 0.34, RED, RED)
    add_oval(slide, x + w * 0.69, y + 3.43, 0.18, 0.18, WHITE, WHITE)
    add_text(slide, x + w * 0.76, y + 3.37, 1.2, 0.12, "対応種別の内訳", 6.3, BLACK, True)


def mock_slide(prs, page: int, title: str, active: str, panel_title: str, bullets: list[str], body: str = "dashboard") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, page, "Chapter2", title)
    ui_x, ui_y, ui_w, ui_h = 0.78, 1.35, 9.35, 5.35
    if body == "dashboard":
        slide.shapes.add_picture(str(MOCK_DASHBOARD_IMAGE), Inches(ui_x), Inches(ui_y), width=Inches(ui_w), height=Inches(ui_h))
        bx, by, bw, bh = ui_x, ui_y, ui_w, ui_h
    elif body == "detail":
        bx, by, bw, bh = app_shell(slide, ui_x, ui_y, ui_w, ui_h, title.replace("スクショ資料：", ""), active)
        add_round(slide, bx, by, 2.25, 1.45, GRAY_CARD, LINE)
        add_text(slide, bx + 0.18, by + 0.18, 1.0, 0.16, "物件情報", 8.5, BLACK, True)
        add_text(slide, bx + 0.18, by + 0.52, 1.75, 0.55, "所在地　東京都〇〇区\n契約状況　更新確認中\n担当者　営業部 A様\n優先度　高", 6.7, TEXT)
        add_rect(slide, bx + 3.05, by + 0.15, 0.02, 3.2, rgb("d1d1d1"), None)
        for i, (h, b) in enumerate([("受付", "問い合わせ登録"), ("確認", "担当者確認"), ("承認", "上長承認"), ("完了", "履歴保存")]):
            yy = by + 0.30 + i * 0.78
            add_oval(slide, bx + 2.92, yy, 0.28, 0.28, RED if i < 3 else rgb("bdbdbd"), None)
            add_text(slide, bx + 3.35, yy - 0.04, 0.65, 0.14, h, 7.8, BLACK, True)
            add_text(slide, bx + 3.35, yy + 0.18, 1.2, 0.12, b, 6.2, MUTED)
        add_round(slide, bx + 5.25, by, 2.35, 3.50, GRAY_CARD, LINE)
        add_text(slide, bx + 5.45, by + 0.18, 1.4, 0.16, "関連書類 / メモ", 8.2, BLACK, True)
        for i, item in enumerate(["契約書ドラフト", "見積依頼", "問い合わせ", "報告メモ", "過去対応履歴"]):
            add_round(slide, bx + 5.45, by + 0.58 + i * 0.48, 1.9, 0.30, WHITE, LINE)
            add_text(slide, bx + 5.58, by + 0.67 + i * 0.48, 1.5, 0.10, item, 5.8, BLACK)
    else:
        bx, by, bw, bh = app_shell(slide, ui_x, ui_y, ui_w, ui_h, title.replace("スクショ資料：", ""), active)
        add_round(slide, bx, by, bw, 0.58, GRAY_CARD, LINE)
        add_text(slide, bx + 0.20, by + 0.20, 2.6, 0.13, "契約更新時の注意点を検索", 7.8, BLACK, True)
        items = [("社内FAQ", "更新案内、請求、インボイス対応の社内手順を提示"), ("契約書", "関連条文と過去ドラフトを候補表示"), ("対応履歴", "類似案件の対応履歴を要約"), ("オーナー報告", "訪問前に物件・収支・過去提案を整理")]
        for i, (h, b) in enumerate(items):
            x = bx + (i % 2) * 3.0
            y = by + 0.90 + (i // 2) * 1.25
            add_round(slide, x, y, 2.55, 0.86, WHITE, LINE)
            add_round(slide, x + 0.18, y + 0.16, 0.88, 0.22, rgb("ffeaea"), rgb("ffeaea"))
            add_text(slide, x + 0.22, y + 0.21, 0.78, 0.09, h, 5.8, RED, True, PP_ALIGN.CENTER)
            add_text(slide, x + 0.22, y + 0.52, 1.95, 0.20, b, 5.6, TEXT)
    add_round(slide, 10.38, 1.45, 2.48, 4.95, WHITE, LINE)
    add_text(slide, 10.68, 1.78, 1.4, 0.30, panel_title, 15, RED, True)
    y = 2.40
    for item in bullets:
        add_oval(slide, 10.68, y + 0.02, 0.08, 0.08, RED, None)
        add_text(slide, 10.84, y - 0.02, 1.65, 0.36, item, 11, TEXT)
        y += 0.72


def roadmap_editable(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 13, "Chapter3", "今後の進め方のイメージ")
    add_text(slide, 0.74, 1.25, 11.8, 0.45, "現場の混乱を避けながら進めるため、下記のような流れで段階的に進める想定です。", 18, BLACK, True)
    add_round(slide, 0.70, 2.05, 12.05, 4.72, WHITE, None)
    steps = [
        ("01", "ご契約", "NDA・個別契約など\n必要な契約を締結"),
        ("02", "進め方の確認", "役割分担を確認\nデータ整理方針を決定"),
        ("03", "移行先の準備", "移行先の箱を準備\n既存データを確認"),
        ("04", "開発", "主要画面から開発\n定例で使用感を確認"),
        ("05", "現場導入", "小規模テスト後に導入\n必要に応じて説明"),
        ("06", "改善", "運用後の改善を反映\n機能追加を検討"),
    ]
    start_x = 1.02
    gap = 1.93
    for i, (num, label, body) in enumerate(steps):
        x = start_x + i * gap
        add_oval(slide, x + 0.22, 2.48, 0.92, 0.92, WHITE, rgb("d1d1d1"))
        add_oval(slide, x + 0.30, 2.56, 0.76, 0.76, None, PINK)
        add_text(slide, x + 0.41, 2.78, 0.54, 0.28, num, 24, RED, True, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_rect(slide, x + 1.18, 2.92, 0.54, 0.02, PINK, PINK)
            add_oval(slide, x + 1.44, 2.86, 0.12, 0.12, RED, None)
        arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, 3.70, 1.55, 0.54, RED, RED)
        add_text(slide, x + 0.05, 3.87, 1.30, 0.12, label, 12.5, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.02, 4.58, 1.55, 1.25, body, 11.0, BLACK)


def comparison(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 14, "Chapter3", "比較検討時の判断軸")
    headers = [("判断軸", 0.82, 2.25), ("一般的なパッケージ導入", 3.05, 3.9), ("Field Xで開発する場合", 7.20, 5.25)]
    for h, x, w in headers:
        add_round(slide, x, 1.55, w, 0.42, RED, RED)
        add_text(slide, x, 1.67, w, 0.14, h, 12.5, WHITE, True, PP_ALIGN.CENTER)
    rows = [
        ("業務適合", "既存機能に業務を寄せる\n必要がある", "貴社の業務フロー・承認・帳票に\n合わせて設計"),
        ("現場定着", "画面や操作が合わず\n利用が限定される可能性", "担当者様の使い方を前提に\nモックから改善"),
        ("AI拡張", "個別AI連携は追加開発・\n制約が出やすい", "基幹データを起点に、音声対応・\n営業支援AIまで展開"),
        ("資料化", "導入効果や運用イメージを\n別途整理する必要", "スクショ資料・要件資料・進め方を\n一式で確認可能"),
    ]
    y = 2.25
    for axis, pkg, fx in rows:
        add_round(slide, 0.82, y, 2.25, 0.78, WHITE, LINE)
        add_round(slide, 3.05, y, 3.9, 0.78, WHITE, LINE)
        add_round(slide, 7.20, y, 5.25, 0.78, rgb("fff1f1"), RED)
        add_text(slide, 1.05, y + 0.24, 1.2, 0.20, axis, 13, BLACK, True)
        add_text(slide, 3.25, y + 0.18, 3.3, 0.36, pkg, 11.5, TEXT)
        add_text(slide, 7.45, y + 0.18, 4.55, 0.36, fx, 11.5, TEXT)
        y += 0.92


def future_ai(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 15, "Chapter3", "基幹管理システムを起点に、社内AIへ展開可能")
    ui_x, ui_y, ui_w, ui_h = 0.85, 1.55, 5.25, 4.25
    bx, by, bw, bh = app_shell(slide, ui_x, ui_y, ui_w, ui_h, "基幹データ", "総合")
    dashboard_body(slide, bx, by, bw, bh)
    add_oval(slide, 6.95, 3.15, 1.55, 0.85, RED, RED)
    add_text(slide, 7.13, 3.34, 1.18, 0.36, "基幹データ\n社内AI基盤", 12.5, WHITE, True, PP_ALIGN.CENTER)
    nodes = [
        (6.70, 1.65, "AI音声対応", "入居者様問い合わせ"),
        (9.55, 1.65, "営業支援AI", "訪問前ブリーフィング"),
        (6.70, 5.10, "契約書自動作成", "文書作成を省力化"),
        (9.55, 5.10, "社内FAQ検索", "横断検索・ナレッジ化"),
    ]
    for x, y, h, b in nodes:
        add_round(slide, x, y, 2.38, 0.78, WHITE, LINE)
        add_text(slide, x + 0.18, y + 0.18, 1.6, 0.18, h, 13, RED, True)
        add_text(slide, x + 0.18, y + 0.48, 1.8, 0.12, b, 9, TEXT)
    add_round(slide, 10.62, 3.05, 1.85, 1.55, WHITE, LINE)
    add_text(slide, 10.88, 3.28, 1.1, 0.20, "展開イメージ", 13, RED, True)
    add_bullets(slide, 10.86, 3.70, 1.35, 0.62, ["問い合わせ一次対応", "訪問前の情報整理", "社内情報をFAQ化"], 8.5)


def build() -> None:
    prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover(prs)
    agenda(prs)
    chapter(prs, 3, "01", "これまでにお伺いしたこと")
    position_slide(prs)
    heard_slide(prs)
    chapter(prs, 6, "02", "ご提案の要旨・デモンストレーション")
    proposal_summary(prs)
    mock_slide(prs, 8, "改善版モックで提示する全体像", "総合", "画面で確認できること", ["左サイドバーを維持", "期限・担当・状態を一覧化", "画面単位で要件を確認"], "dashboard")
    mock_slide(prs, 9, "スクショ資料：基幹管理ダッシュボード", "総合", "主な機能", ["本日期限と遅延リスクを可視化", "担当者別・案件種別で確認", "対応履歴と関連書類を連携"], "dashboard")
    mock_slide(prs, 10, "スクショ資料：物件・契約ごとの進捗管理", "物件", "主な機能", ["物件・契約・履歴を一画面管理", "承認から完了まで可視化", "契約書自動作成へ接続"], "detail")
    mock_slide(prs, 11, "スクショ資料：社内FAQ・横断検索", "検索", "拡張時の価値", ["契約書・規程・履歴を横断検索", "担当者ごとの差を小さくする", "問い合わせ・営業準備に展開"], "search")
    chapter(prs, 12, "03", "今後の進め方のイメージ")
    roadmap_editable(prs)
    comparison(prs)
    future_ai(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    with zipfile.ZipFile(OUT) as z:
        bad = z.testzip()
        if bad:
            raise RuntimeError(f"Broken pptx entry: {bad}")
    print(OUT)


if __name__ == "__main__":
    build()
